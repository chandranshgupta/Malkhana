import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme colors
    c_dark = RGBColor(15, 23, 42)      # #0F172A (Ink)
    c_blue = RGBColor(14, 165, 233)    # #0EA5E9 (Accent Blue)
    c_light = RGBColor(244, 247, 249)  # #F4F7F9 (Paper background)
    c_white = RGBColor(255, 255, 255)
    c_gray = RGBColor(100, 116, 139)   # #64748B (Muted text)
    c_red = RGBColor(220, 38, 38)      # #DC2626 (Alert/Red)
    c_green = RGBColor(22, 163, 74)    # #16A34A (OK Green)
    
    blank_layout = prs.slide_layouts[6]
    
    # Helper to add full background
    def set_slide_background(slide, color):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.color.rgb = color
        # Send to back is done by placing it first on the slide
        return rect

    # Helper to add blueprint corner brackets
    def add_corner_marks(slide, color):
        # Top-left corner
        # We can draw small L-shapes or just simple lines
        # Instead of complex drawing, we can add thin border rectangles or text marks
        pass

    # ==================== SLIDE 1: Cover (Dark theme) ====================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, c_dark)
    
    # Title Text Frame
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "MALKHANA VAULT"
    p0.font.name = "Calibri"
    p0.font.size = Pt(54)
    p0.font.bold = True
    p0.font.color.rgb = c_blue
    p0.alignment = PP_ALIGN.LEFT
    
    p1 = tf.add_paragraph()
    p1.text = "Digital Chain of Custody for Indian Law Enforcement"
    p1.font.name = "Calibri"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = c_white
    p1.space_before = Pt(12)
    p1.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = "BSA 2023 §63 & BNSS 2023 §153 Compliant | Offline-First Forensic-Grade System"
    p2.font.name = "Courier New"
    p2.font.size = Pt(13)
    p2.font.color.rgb = c_gray
    p2.space_before = Pt(36)
    p2.alignment = PP_ALIGN.LEFT

    # ==================== SLIDE 2: The Problem (Light theme) ====================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, c_light)
    
    # Title
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "THE PROBLEM: MANUAL MALKHANA RISKS"
    p.font.name = "Calibri"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = c_dark
    
    # Subtitle / Paragraphs
    tb_content = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0))
    tf_c = tb_content.text_frame
    tf_c.word_wrap = True
    
    p_intro = tf_c.paragraphs[0]
    p_intro.text = "Physical property rooms ('Malkhanas') suffer from critical gaps that lead to evidence rejection in court:"
    p_intro.font.name = "Calibri"
    p_intro.font.size = Pt(18)
    p_intro.font.color.rgb = c_dark
    p_intro.space_after = Pt(20)
    
    problems = [
        ("No Crime Scene Hash (H1 Gap):", "Devices are seized without immediate hash values, allowing defense attorneys to challenge baseline integrity."),
        ("Physical Chain Alterations:", "Manual registers are vulnerable to backdating, missing entries, and paper logs can be tampered with or damaged."),
        ("No Write-Blocker Enforcement:", "Evidence is often plugged directly into station computers without write-blockers, modifying file timestamps and metadata."),
        ("Inadmissible Certifications:", "Section 65B (IEA) or Section 63 (BSA) certificates are drafted manually, leading to technical errors and court rejection.")
    ]
    
    for title, desc in problems:
        p_item = tf_c.add_paragraph()
        p_item.space_before = Pt(10)
        p_item.space_after = Pt(10)
        
        # Add bold bullet title
        run_title = p_item.add_run()
        run_title.text = "• " + title + " "
        run_title.font.name = "Calibri"
        run_title.font.bold = True
        run_title.font.size = Pt(16)
        run_title.font.color.rgb = c_red
        
        # Add description
        run_desc = p_item.add_run()
        run_desc.text = desc
        run_desc.font.name = "Calibri"
        run_desc.font.size = Pt(16)
        run_desc.font.color.rgb = c_dark

    # ==================== SLIDE 3: The Solution (Light theme) ====================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, c_light)
    
    # Title
    tb = slide3.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "THE SOLUTION: MALKHANA VAULT"
    p.font.name = "Calibri"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = c_dark
    
    # Content Columns (Two-column layout)
    # Left Column - Highlights
    tb_left = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.5))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    
    p_l = tf_l.paragraphs[0]
    p_l.text = "Forensic-Grade Air-Gapped Trust"
    p_l.font.name = "Calibri"
    p_l.font.bold = True
    p_l.font.size = Pt(20)
    p_l.font.color.rgb = c_blue
    p_l.space_after = Pt(12)
    
    features = [
        "100% Offline-First execution, eliminating cloud hacks.",
        "SQLCipher database with AES-256 key encryption.",
        "Append-only database architecture—no DELETE queries compile.",
        "Secure auth portal records session as Step 0 custody."
    ]
    for feat in features:
        p_f = tf_l.add_paragraph()
        p_f.text = "✔  " + feat
        p_f.font.name = "Calibri"
        p_f.font.size = Pt(15)
        p_f.font.color.rgb = c_dark
        p_f.space_before = Pt(8)
        
    # Right Column - Visual representation (Simulated Terminal Box)
    rect_term = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5))
    rect_term.fill.solid()
    rect_term.fill.fore_color.rgb = c_dark
    rect_term.line.color.rgb = c_blue
    rect_term.line.width = Pt(1.5)
    
    tf_t = rect_term.text_frame
    tf_t.margin_top = Inches(0.3)
    tf_t.margin_left = Inches(0.3)
    tf_t.word_wrap = True
    
    t_lines = [
        ("$ chmod +x MalkhanaVault.AppImage", c_blue, True),
        ("$ ./MalkhanaVault.AppImage", c_blue, True),
        ("[*] Initializing SQLCipher database...", c_gray, False),
        ("[*] Verifying Merkle audit chain...", c_gray, False),
        ("✔ Vault unlocked. Chain-of-custody active.", c_green, True),
        ("$ malkhana --verify-chain", c_blue, True),
        ("✔ Merkle root verified. 0 tamper events.", c_green, True)
    ]
    
    for i, (text, color, bold) in enumerate(t_lines):
        if i == 0:
            p_t = tf_t.paragraphs[0]
        else:
            p_t = tf_t.add_paragraph()
        p_t.text = text
        p_t.font.name = "Courier New"
        p_t.font.size = Pt(11)
        p_t.font.color.rgb = color
        p_t.font.bold = bold
        p_t.space_before = Pt(6)

    # ==================== SLIDE 4: The 150-Drawer Vault Matrix (Light theme) ====================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, c_light)
    
    # Title
    tb = slide4.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "THE 150-DRAWER SECURE VAULT MATRIX"
    p.font.name = "Calibri"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = c_dark
    
    # Text explanation on the left
    tb_exp = slide4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_exp = tb_exp.text_frame
    tf_exp.word_wrap = True
    
    p_exp_t = tf_exp.paragraphs[0]
    p_exp_t.text = "Physical Storage Matrix Alignment"
    p_exp_t.font.name = "Calibri"
    p_exp_t.font.bold = True
    p_exp_t.font.size = Pt(20)
    p_exp_t.font.color.rgb = c_blue
    p_exp_t.space_after = Pt(12)
    
    exp_bullets = [
        "Coordinate Mapping: Bridges physical evidence drawers to database slots (e.g. Row 5, Column 8).",
        "Dual Verification: Enforces H2 Receipt Hash matches crime scene H1 before drawer is digitally sealed.",
        "Malkhana Gap Alert: Automatic system warning if H1 does not match H2 on receipt, preventing legal fallout."
    ]
    for bullet in exp_bullets:
        p_b = tf_exp.add_paragraph()
        p_b.text = "• " + bullet
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(15)
        p_b.font.color.rgb = c_dark
        p_b.space_before = Pt(10)
        
    # Draw a visual grid on the right (simulate the 150-drawer grid)
    # Draw a 10x6 small grid of squares to represent 60 of the drawers visually
    grid_left = 6.0
    grid_top = 1.8
    cell_size = 0.55
    gap = 0.08
    
    for row in range(6):
        for col in range(10):
            l = grid_left + col * (cell_size + gap)
            t = grid_top + row * (cell_size + gap)
            cell = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(cell_size), Inches(cell_size))
            
            # Highlight a few drawers (Row 3 Col 4, Row 1 Col 8 etc.)
            if (row == 2 and col == 3) or (row == 0 and col == 7):
                cell.fill.solid()
                cell.fill.fore_color.rgb = c_blue # Occupied/Sealed
                cell.line.color.rgb = c_blue
            elif (row == 4 and col == 5):
                cell.fill.solid()
                cell.fill.fore_color.rgb = c_red # Alert mismatch
                cell.line.color.rgb = c_red
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = c_white # Empty slot
                cell.line.color.rgb = c_gray
                cell.line.width = Pt(0.5)

    # Add legends below the grid
    tb_leg = slide4.shapes.add_textbox(Inches(6.0), Inches(5.8), Inches(6.3), Inches(1.0))
    tf_leg = tb_leg.text_frame
    p_leg = tf_leg.paragraphs[0]
    p_leg.text = "Matrix Legends:  [ ] Empty Slot      [■] Occupied/Sealed      [■] Integrity Alert"
    p_leg.font.name = "Courier New"
    p_leg.font.size = Pt(11)
    p_leg.font.bold = True
    p_leg.font.color.rgb = c_dark
    p_leg.font.color.rgb = c_gray

    # ==================== SLIDE 5: Legal Framework (Light theme) ====================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, c_light)
    
    # Title
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.3), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "BSA 2023 & BNSS 2023 COMPLIANCE"
    p.font.name = "Calibri"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = c_dark
    
    # Legal cells (simulating the 2x2 grid from website)
    cell_data = [
        ("BSA 2023 — §63", "Section 63 Certificate", "Auto-generated certificate attesting device health, automated storage, integrity, and chain of custody. Replaces legacy IEA §65B."),
        ("BSA 2023 — §63(2)(c)", "Downtime Logging", "System interruptions are auto-flagged in system_health_log. Certificate generation checks and discloses any downtime events."),
        ("BNSS 2023 — §153", "Form CC-1 Integration", "Mandatory chain-of-custody form auto-populated at seizure. Enforces Panch witness names and Imaging Specialist signatures."),
        ("CERT-In / IT Act", "180-Day Logs", "Tamper-proof Merkle-chained audit logs preserved for a minimum of 180 days. Any stealth edit breaks the Merkle root validation.")
    ]
    
    col_w = 5.2
    row_h = 2.1
    c_positions = [
        (1.0, 1.8), # Cell 1 (Left top)
        (6.8, 1.8), # Cell 2 (Right top)
        (1.0, 4.2), # Cell 3 (Left bottom)
        (6.8, 4.2)  # Cell 4 (Right bottom)
    ]
    
    for i, (tag, heading, desc) in enumerate(cell_data):
        l, t = c_positions[i]
        cell_rect = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(col_w), Inches(row_h))
        cell_rect.fill.solid()
        cell_rect.fill.fore_color.rgb = c_white
        cell_rect.line.color.rgb = c_gray
        cell_rect.line.width = Pt(0.5)
        
        tf_cell = cell_rect.text_frame
        tf_cell.margin_top = Inches(0.2)
        tf_cell.margin_left = Inches(0.2)
        tf_cell.word_wrap = True
        
        p_tag = tf_cell.paragraphs[0]
        p_tag.text = tag
        p_tag.font.name = "Courier New"
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = c_blue
        
        p_head = tf_cell.add_paragraph()
        p_head.text = heading
        p_head.font.name = "Calibri"
        p_head.font.size = Pt(16)
        p_head.font.bold = True
        p_head.font.color.rgb = c_dark
        p_head.space_before = Pt(4)
        
        p_desc = tf_cell.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = c_gray
        p_desc.space_before = Pt(6)

    # ==================== SLIDE 6: Summary & Deployment (Dark theme) ====================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, c_dark)
    
    # Title
    tb = slide6.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "MALKHANA VAULT: SECURE & DEPLOYABLE"
    p.font.name = "Calibri"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = c_white
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = "Standardized packages designed for grassroots police deployment and air-gapped security."
    p_sub.font.name = "Calibri"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = c_blue
    p_sub.space_before = Pt(8)
    
    # Four deployments text
    tb_tiers = slide6.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.3), Inches(3.0))
    tf_t = tb_tiers.text_frame
    tf_t.word_wrap = True
    
    deployments = [
        ("AppImage (Tauri Portable)", "Runs instantly on Linux workstations without installation or administrator privileges."),
        ("Debian / RPM Installers", "Native packages for official state Linux distributions (Ubuntu, BOSS Linux)."),
        ("Windows MSIs (NSIS)", "Standard installers with automatic dependency resolution for Windows terminals."),
        ("100% Offline Integrity", "Statically compiled SQLite/SQLCipher database. Zero external web calls or data leaks.")
    ]
    
    for i, (title, desc) in enumerate(deployments):
        if i == 0:
            p_ti = tf_t.paragraphs[0]
        else:
            p_ti = tf_t.add_paragraph()
        p_ti.space_before = Pt(8)
        
        r_t = p_ti.add_run()
        r_t.text = "■  " + title + ": "
        r_t.font.name = "Calibri"
        r_t.font.bold = True
        r_t.font.size = Pt(15)
        r_t.font.color.rgb = c_blue
        
        r_d = p_ti.add_run()
        r_d.text = desc
        r_d.font.name = "Calibri"
        r_d.font.size = Pt(15)
        r_d.font.color.rgb = c_white
        
    # Save presentation
    os.makedirs("collateral", exist_ok=True)
    prs.save("collateral/pitch-deck.pptx")
    print("Presentation saved successfully at collateral/pitch-deck.pptx")

if __name__ == "__main__":
    create_presentation()
